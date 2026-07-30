"""Builds AOS_Unicity_Investment_Deck.pptx (12 slides). Converted to PDF separately."""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))
import facts as F
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY=RGBColor(0x1F,0x2A,0x44); GOLD=RGBColor(0xB8,0x86,0x0B); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREY=RGBColor(0x5A,0x5F,0x68); LGREY=RGBColor(0xF2,0xF3,0xF5); RED=RGBColor(0xA0,0x30,0x30)
DKG=RGBColor(0x2E,0x2E,0x2E)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,fill,line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line
    sh.shadow.inherit=False
    return sh
def txt(s,x,y,w,h,text,size=14,color=DKG,bold=False,align=PP_ALIGN.LEFT,font="Calibri",anchor=MSO_ANCHOR.TOP,line_sp=1.05):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2);tf.margin_right=Pt(2);tf.margin_top=Pt(1);tf.margin_bottom=Pt(1)
    lines = text if isinstance(text,list) else [text]
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line_sp
        if isinstance(ln,tuple): s_,bl_ = ln
        else: s_,bl_ = ln, bold
        r=p.add_run(); r.text=s_; f=r.font; f.size=Pt(size); f.name=font; f.bold=bl_; f.color.rgb=color
    return tb
def bullets(s,x,y,w,h,items,size=13,color=DKG,gap=4,bullet="•  "):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.line_spacing=1.03
        lvl0 = it[0] if isinstance(it,tuple) else it
        clr = it[1] if isinstance(it,tuple) else color
        r=p.add_run(); r.text=bullet+lvl0; r.font.size=Pt(size); r.font.name="Calibri"; r.font.color.rgb=clr
    return tb
def header(s,kicker,title):
    rect(s,0,0,SW,Inches(1.15),NAVY)
    rect(s,0,Inches(1.15),SW,Pt(3),GOLD)
    txt(s,Inches(0.5),Inches(0.12),Inches(12),Inches(0.3),kicker,11,GOLD,True,font="Calibri")
    txt(s,Inches(0.5),Inches(0.4),Inches(12.3),Inches(0.7),title,25,WHITE,True)
def footer(s,n):
    txt(s,Inches(0.5),Inches(7.08),Inches(9),Inches(0.3),"Unicity Labs / AOS Thesis  |  Sahil Modi for Headline  |  July 2026  |  USD",8,GREY)
    txt(s,Inches(12.2),Inches(7.08),Inches(0.9),Inches(0.3),f"{n} / 12",8,GREY,align=PP_ALIGN.RIGHT)

def chip(s,x,y,label,val,w=Inches(1.75)):
    rect(s,x,y,w,Inches(0.95),LGREY)
    txt(s,x,y+Inches(0.12),w,Inches(0.45),val,20,NAVY,True,PP_ALIGN.CENTER)
    txt(s,x,y+Inches(0.58),w,Inches(0.32),label,9.5,GREY,False,PP_ALIGN.CENTER)

b=F.BROAD
# ---------------- Slide 1: Recommendation + funnel (lead with AOS) ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE)
rect(s,0,0,SW,Inches(1.7),NAVY); rect(s,0,Inches(1.7),SW,Pt(4),GOLD)
txt(s,Inches(0.5),Inches(0.22),Inches(12),Inches(0.35),"HEADLINE TAKE-HOME  |  SELECTED OPPORTUNITY",12,GOLD,True)
txt(s,Inches(0.5),Inches(0.46),Inches(12.4),Inches(0.9),"AOS could become the control and proof layer beneath autonomous agents",23,WHITE,True,line_sp=1.0)
txt(s,Inches(0.5),Inches(1.42),Inches(12.4),Inches(0.35),"AOS / Unicity Labs. Selected company for focused diligence. Interesting investment opportunity, not an immediate check.",12.5,RGBColor(0xD8,0xDD,0xE6))
txt(s,Inches(0.5),Inches(2.0),Inches(12.3),Inches(0.9),
    [("Advance into focused diligence. Do not write a check now.",True),
     ("The opportunity has foundational infrastructure potential, but product status, enterprise demand, equity value capture, and financing remain unresolved.",False)],
    14,DKG)
# funnel chips
y=Inches(3.15)
for i,(lab,val) in enumerate([("Posts returned",f"{b['returned']:,}"),("Net-new Posts",f"{b['net_new']:,}"),
     ("Actionable",f"{b['actionable']}"),("Companies",f"{b['consolidated']}"),("Profiles enriched",f"{b['enriched']}"),
     ("Est. activity",f"${F.COST_TOTAL}")]):
    chip(s,Inches(0.5+i*2.06),y,lab,val)
txt(s,Inches(0.5),Inches(4.35),Inches(12.3),Inches(0.3),"Broad-market run. Estimated API activity of the 25 dollar allowance. The engine searched X's recent seven-day window with approved queries; it did not review the entire market.",10,GREY)
# honesty box
rect(s,Inches(0.5),Inches(4.85),Inches(12.33),Inches(1.75),LGREY)
rect(s,Inches(0.5),Inches(4.85),Pt(5),Inches(1.75),GOLD)
txt(s,Inches(0.75),Inches(5.0),Inches(11.9),Inches(1.6),
    [("Honest framing",True),
     ("AOS did not win the standardized scorecard. ScaleDown scored higher (5.95 versus 5.55) with clearer current product documentation and financing actionability. Vattara had simpler pre-seed timing and a live SaaS model. I selected AOS because the scorecard is a decision aid, not the decision. AOS has the strongest team signal, the deepest architecture, and the most asymmetric platform outcome, and I accept its higher Web3 and value-capture risk.",False)],
    12.5,DKG)
footer(s,1)

# ---------------- Slide 2: Why agents require runtime governance ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"THE PROBLEM AND THE OPPORTUNITY","Why this is an interesting investment")
txt(s,Inches(0.5),Inches(1.32),Inches(6.1),Inches(0.4),"The control problem",14,NAVY,True)
bullets(s,Inches(0.5),Inches(1.82),Inches(6.1),Inches(4.9),[
 "Agents are shifting from generating content to taking actions: calling tools, moving value, coordinating.",
 "Prompt-level guardrails do not bind a model that can be talked out of its instructions.",
 "Once agents hold economic authority, enterprises need controls the model cannot override.",
 "Those controls belong in the execution path, below the model: identity, permissions, policy, budgets, isolation, settlement, and audit.",
],13,gap=8)
rect(s,Inches(6.85),Inches(1.3),Inches(5.98),Inches(5.15),LGREY)
txt(s,Inches(7.1),Inches(1.45),Inches(5.5),Inches(0.4),"Why this is an interesting investment",14,NAVY,True)
bullets(s,Inches(7.1),Inches(1.95),Inches(5.5),Inches(4.4),[
 "Foundational control problem: enforcement the model cannot bypass.",
 "Horizontal infrastructure: beneath many agent applications, not one workflow.",
 "Strong technical-team signal: cryptography, distributed systems, Guardtime (company-reported).",
 "Growing agent authority: the need scales as agents transact.",
 "Potential cryptographic differentiation: proof at the execution boundary, not a prompt filter.",
 "Asymmetric platform outcome: part of the control plane for autonomous software.",
],12.5,gap=7)
txt(s,Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.3),"Headline has publicly named agent authentication, runtimes and sandboxes, MCP integration, and orchestration as priority infrastructure problems. This does not imply Headline has expressed interest in Unicity.",10,GREY)
footer(s,2)

# ---------------- Slide 3: AOS and Unicity entity architecture ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"ENTITY AND PRODUCT MAP","AOS and the Unicity entity architecture")
txt(s,Inches(0.5),Inches(1.35),Inches(12.3),Inches(0.4),"These are not interchangeable. Distinguishing them is diligence question one.",12,RED,True)
rows=[("Unicity Labs","Operating company, Zug. Team and presumed equity value. IP ownership unresolved.","company-reported"),
 ("Unicity Foundation","Swiss foundation. Protocol governance, grants, open source. May hold economics.","company-reported"),
 ("Unicity Protocol","Bearer-token, peer-to-peer settlement. The live product today. Web3 infrastructure.","official source"),
 ("AOS (working thesis)","Secure agent execution, identity, policy, budget, proof layer. Not a confirmed current brand.","analyst hypothesis"),
 ("Agent Sphere / Sphere SDK","Agent wallet and payments surface. Marketplace mechanics unconfirmed.","official source"),
 ("Bearer tokens","Self-contained, self-proving assets transferable peer to peer. Money for machines.","official source")]
y=Inches(1.85); rh=Inches(0.78)
for i,(a,d,c) in enumerate(rows):
    bg=LGREY if i%2 else WHITE
    rect(s,Inches(0.5),y,Inches(12.33),rh,bg)
    txt(s,Inches(0.65),y+Inches(0.06),Inches(2.6),rh,a,13,NAVY,True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.35),y+Inches(0.06),Inches(7.3),rh,d,12,DKG,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(10.8),y+Inches(0.06),Inches(1.95),rh,c,10.5,GREY,anchor=MSO_ANCHOR.MIDDLE)
    y=y+rh+Inches(0.02)
txt(s,Inches(0.5),Inches(6.9),Inches(12.3),Inches(0.3),"Closest current artifacts to an agent execution layer: an Unicity Execution Model specification and a tx-flow runtime in a 79-repo GitHub organization. No repo named AOS or Astrid in the authoritative listing.",10,GREY)
footer(s,3)

# ---------------- Slide 4: Team, evidence, funding, product status ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"TEAM, EVIDENCE, FUNDING, STATUS","The strongest signal, and the honest gaps")
txt(s,Inches(0.5),Inches(1.4),Inches(6),Inches(0.4),"Team (company-reported)",14,NAVY,True)
bullets(s,Inches(0.5),Inches(1.85),Inches(6.1),Inches(2.1),[
 "Mike Gault, CEO, previously built and exited Guardtime, a cryptographic infrastructure company.",
 "Team reported to include PhD researchers in distributed systems, cryptography, and machine learning.",
 "This is the strongest team signal in the finalist set.",
],13,gap=7)
txt(s,Inches(0.5),Inches(4.0),Inches(6),Inches(0.4),"Funding (company-reported PR)",14,NAVY,True)
bullets(s,Inches(0.5),Inches(4.45),Inches(6.1),Inches(2.1),[
 "3 million dollar seed announced February 19, 2026.",
 "Led by Blockchange Ventures; Outlier Ventures and Tawasal participating.",
 "Unicity Foundation established in Switzerland.",
 "The round is closed, which reduces immediate financing actionability.",
],13,gap=6)
txt(s,Inches(6.9),Inches(1.4),Inches(6),Inches(0.4),"Current product status (verified direction)",14,NAVY,True)
bullets(s,Inches(6.9),Inches(1.85),Inches(5.9),Inches(2.3),[
 "The live site presents bearer tokens and peer-to-peer settlement, with a developer SDK named Sphere.",
 "Performance claims of 3M+ TPS and under one second finality are company-reported, not reproduced.",
 "A distinct enterprise AOS product is unresolved.",
],13,gap=7)
txt(s,Inches(6.9),Inches(4.0),Inches(6),Inches(0.4),"Traction (unresolved)",14,NAVY,True)
bullets(s,Inches(6.9),Inches(4.45),Inches(5.9),Inches(2.1),[
 ("No independently verified paid customer or production deployment exists in the saved diligence.",RED),
 "GitHub activity across 79 repositories is not commercial traction.",
 "Funding does not prove enterprise demand.",
],13,gap=7)
footer(s,4)

# ---------------- Slide 5: Competition, business model, value capture ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"COMPETITION AND VALUE CAPTURE","Landscape, business model, and the equity question")
txt(s,Inches(0.5),Inches(1.35),Inches(6.2),Inches(0.4),"Competitive layers (illustrative, no invented data)",13,NAVY,True)
bullets(s,Inches(0.5),Inches(1.8),Inches(6.2),Inches(4.6),[
 "Runtime and sandbox: E2B, Modal, Daytona, Fly.io.",
 "Identity: Auth0, Stytch, Descope, WorkOS, MCP auth.",
 "Governance and security: Lakera, Protect AI, HiddenLayer, Prompt Security, CalypsoAI.",
 "Payments: Skyfire, Coinbase x402, Google AP2, Nevermined, stablecoin rails.",
 "Protocols: MCP, A2A, Fetch.ai.",
 "Internal build: cloud IAM, service meshes, API gateways, audit logging.",
],12.5,gap=7)
txt(s,Inches(6.9),Inches(1.35),Inches(6),Inches(0.4),"Value-capture paths (all hypotheses)",13,NAVY,True)
bullets(s,Inches(6.9),Inches(1.8),Inches(5.9),Inches(2.6),[
 "Enterprise license: annual platform fee for governance, runtime, policy, identity, audit.",
 "Usage-based: fee per action, execution, proof, or settlement.",
 "Developer platform: free tier, paid production, enterprise controls.",
 "Protocol and token economics: network fees, token value, foundation-controlled.",
],12.5,gap=6)
rect(s,Inches(6.9),Inches(4.55),Inches(5.9),Inches(2.0),LGREY); rect(s,Inches(6.9),Inches(4.55),Pt(5),Inches(2.0),RED)
txt(s,Inches(7.1),Inches(4.7),Inches(5.5),Inches(1.75),
 [("The central value-capture question",True),
  ("If Unicity becomes successful, which entity receives the economic value, and how does that value reach an equity investor in Unicity Labs rather than a token or the Foundation?",False)],12.5,DKG)
footer(s,5)

# ---------------- Slide 6: Recommendation, risks, diligence ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"RECOMMENDATION","Posture, top risks, and the diligence plan")
rect(s,Inches(0.5),Inches(1.35),Inches(12.33),Inches(1.05),NAVY)
txt(s,Inches(0.75),Inches(1.45),Inches(11.9),Inches(0.9),
 "Advance Unicity Labs into focused diligence. Do not write a check until the four diligence gates are resolved. The announced seed is not presented as an open transaction; this is relationship development and next-round positioning.",
 13.5,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.5),Inches(2.6),Inches(6),Inches(0.4),"Top risks",14,RED,True)
bullets(s,Inches(0.5),Inches(3.05),Inches(6.2),Inches(3.7),[
 "AOS may not be a distinct active product.",
 "Positioning centers on bearer finance, which is Web3 infrastructure.",
 "Token or protocol value may not accrue to equity.",
 "Enterprises may resist protocol complexity.",
 "Cloud, agent frameworks, security vendors, and payment protocols may absorb it.",
 "Seed already announced; funding does not prove demand.",
 "No verified paid customer in the saved diligence.",
],12,gap=4)
txt(s,Inches(6.9),Inches(2.6),Inches(6),Inches(0.4),"The four diligence gates",14,NAVY,True)
bullets(s,Inches(6.9),Inches(3.05),Inches(5.9),Inches(3.7),[
 "1) Product truth: what AOS is today and how it relates to the Execution Model, tx-flow runtime, Sphere, and bearer tokens.",
 "2) Enterprise demand: paid customers, pilots, the economic buyer, and willingness to pay.",
 "3) Equity value capture: which entity owns the IP and how protocol success reaches Labs equity.",
 "4) Financing availability: capitalization, next-round timing, and whether Headline could participate.",
],12,gap=8)
txt(s,Inches(0.5),Inches(6.9),Inches(12.3),Inches(0.3),"Standardized scorecard favored ScaleDown (5.95 vs 5.55). AOS selected as the most interesting company for diligence, not the highest-evidence one. High potential, high uncertainty.",10,GREY)
footer(s,6)

# ---------------- Slide 7: Engine methodology (appendix) ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","Full engine methodology")
bullets(s,Inches(0.5),Inches(1.45),Inches(12.2),Inches(5.4),[
 "Official X API only: search /2/tweets/search/recent, counts /2/tweets/counts/recent, users /2/users. No scraping.",
 "Deterministic and rule-based. No LLM in the classification or disposition path. Same input and config produce the same result.",
 "Attribution separated direct builders from third-party announcements, commentary, and established-company releases.",
 "Level A artifacts (repos, sites, docs, API pages, demos, changelogs, specs) proved existence, never ownership, formation, traction, or revenue.",
 "Artifact resolution order: unwound_url, then expanded_url, then original. X-only links excluded.",
 "Cross-run and within-run deduplication with full query and page provenance.",
 "Consolidation on exact repo, domain, or project plus author. Duplicate evidence did not raise scores.",
 "Every paid request was fingerprinted with SHA-256, typed-approved, budget-capped in exact decimals, and locked one-time against re-execution. Raw responses stored before any derived output.",
 "450 passing tests after the final offline enrichment re-derivation.",
],12.5,gap=6)
footer(s,7)

# ---------------- Slide 8: Query and sector coverage ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","Query and sector coverage")
txt(s,Inches(0.5),Inches(1.35),Inches(12),Inches(0.4),"Pilot: 6 query families across 3 lanes. Broad: 20 query families across software, hardware, and physical products.",12,DKG)
groups=[("AI infrastructure and software","q01 to q05"),("Non-AI B2B software","q06 to q13"),
 ("Cybersecurity and fintech","within q06 to q13"),("Robotics and semiconductors","q14, q15"),
 ("Climate and energy hardware","q16"),("Medical devices","q17"),
 ("Industrial and advanced manufacturing","q13, q18"),("Consumer hardware","q19"),("Cross-sector founder transition","q20")]
y=Inches(1.9)
for i,(g,r) in enumerate(groups):
    col = i//5; row=i%5
    x=Inches(0.5+col*6.3); yy=Inches(1.9+row*0.95)
    rect(s,x,yy,Inches(6.0),Inches(0.8),LGREY)
    txt(s,x+Inches(0.2),yy+Inches(0.08),Inches(5.6),Inches(0.35),g,13,NAVY,True)
    txt(s,x+Inches(0.2),yy+Inches(0.45),Inches(5.6),Inches(0.3),r,11,GREY)
txt(s,Inches(0.5),Inches(6.8),Inches(12.3),Inches(0.3),"Exact literal query text for all 6 pilot and 20 broad queries is reproduced verbatim in the Engine Study Guide.",10,GREY)
footer(s,8)

# ---------------- Slide 9: Candidate comparison ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","Candidate comparison, including ScaleDown and Vattara")
txt(s,Inches(0.5),Inches(1.35),Inches(12.3),Inches(0.4),"Standardized scorecard totals (corrected). AOS was NOT the top score. Selection was human judgment.",12,RED,True)
data=[("Company","Standardized total","Note")]
rows=[("ScaleDown","5.95","Higher score: clearer product docs, financing actionability, conventional AI infra"),
 ("AOS / Unicity Labs","5.55","Selected: strongest team, deepest architecture, largest asymmetric upside"),
 ("Vattara","5.20","Live SaaS, public pricing, but crowded and sub-scale raise"),
 ("Verifyr","3.10","Unverified entity"),
 ("Plexor","3.00","Engineering services, not scalable software")]
y=Inches(1.95)
rect(s,Inches(0.5),y,Inches(12.33),Inches(0.5),NAVY)
txt(s,Inches(0.65),y+Inches(0.06),Inches(3),Inches(0.4),"Company",12,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(3.9),y+Inches(0.06),Inches(2.2),Inches(0.4),"Standardized total",12,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(6.2),y+Inches(0.06),Inches(6.4),Inches(0.4),"Note",12,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
y=y+Inches(0.5)
for i,(c,t,n) in enumerate(rows):
    bg=RGBColor(0xED,0xE7,0xD4) if c.startswith("AOS") else (LGREY if i%2 else WHITE)
    rect(s,Inches(0.5),y,Inches(12.33),Inches(0.7),bg)
    txt(s,Inches(0.65),y+Inches(0.05),Inches(3.1),Inches(0.6),c,12.5,NAVY,c.startswith("AOS"),anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.9),y+Inches(0.05),Inches(2.2),Inches(0.6),t,14,DKG,True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(6.2),y+Inches(0.05),Inches(6.4),Inches(0.6),n,11,DKG,anchor=MSO_ANCHOR.MIDDLE)
    y=y+Inches(0.72)
txt(s,Inches(0.5),Inches(6.7),Inches(12.3),Inches(0.4),"The workbook recomputes weighted scores live from editable weights on the Inputs sheet; the standardized totals above are the frozen reference.",10,GREY)
footer(s,9)

# ---------------- Slide 10: Business model and unit economics scenarios ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","Business-model and unit-economics scenarios")
txt(s,Inches(0.5),Inches(1.35),Inches(12.3),Inches(0.4),"Illustrative diligence framework, not a company forecast. All figures are hypotheses until pricing or contracts are verified.",12,RED,True)
bullets(s,Inches(0.5),Inches(1.95),Inches(6.1),Inches(4.6),[
 "Enterprise SaaS: annual platform fee for governance, runtime, policy, identity, and audit.",
 "Usage-based: fee per agent action, execution, proof, or settlement event.",
 "Developer platform: free tier, paid production usage, enterprise controls.",
 "Protocol economics: network or settlement fees, with a possible mismatch between protocol success and equity value.",
],13,gap=8)
txt(s,Inches(6.9),Inches(1.95),Inches(6),Inches(0.4),"Unit economics: three scenarios",13,NAVY,True)
bullets(s,Inches(6.9),Inches(2.45),Inches(5.9),Inches(4.0),[
 "Conservative: few enterprise customers, lower license and usage, thinner margin.",
 "Base: illustrative 12 customers, mid license plus usage, 80 percent gross margin.",
 "Upside: more customers, higher usage, stronger margin and revenue per employee.",
 "The workbook computes revenue, gross profit, contribution, operating profit, revenue per employee, and burn multiple by formula.",
],13,gap=8)
footer(s,10)

# ---------------- Slide 11: AOS vs Protocol vs Foundation vs Agent Sphere ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","AOS versus Protocol versus Foundation versus Agent Sphere")
cards=[("Unicity Protocol","Live product. Bearer-token peer-to-peer settlement. Web3 infrastructure. Consensus layer, state-transition SDKs, yellowpaper.",NAVY),
 ("AOS (thesis)","The investment thesis: secure agent execution and proof layer. Not a confirmed current brand. Diligence question one.",GOLD),
 ("Unicity Foundation","Swiss foundation. Protocol governance, grants, open source. May hold protocol economics separate from company equity.",GREY),
 ("Agent Sphere / Sphere","Agent wallet and payments SDK. Developer and agent surface. Marketplace mechanics unconfirmed.",RGBColor(0x4A,0x5A,0x74))]
for i,(t,d,c) in enumerate(cards):
    x=Inches(0.5+(i%2)*6.3); y=Inches(1.6+(i//2)*2.6)
    rect(s,x,y,Inches(6.0),Inches(2.35),LGREY); rect(s,x,y,Inches(6.0),Inches(0.55),c)
    txt(s,x+Inches(0.2),y+Inches(0.1),Inches(5.6),Inches(0.4),t,15,WHITE,True)
    txt(s,x+Inches(0.2),y+Inches(0.7),Inches(5.6),Inches(1.5),d,12.5,DKG)
txt(s,Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.3),"Public materials sometimes blur these. Separating them, and confirming which one an equity investor actually buys, is the core of first-call diligence.",10,GREY)
footer(s,11)

# ---------------- Slide 12: Source map and evidence discipline ----------------
s=slide(); rect(s,0,0,SW,SH,WHITE); header(s,"APPENDIX","Source map and evidence discipline")
bullets(s,Inches(0.5),Inches(1.45),Inches(6.1),Inches(5.2),[
 "Unicity website and Agent Sphere: official company source.",
 "GitHub organization, 79 repos: official protocol source. Activity is not traction.",
 "3 million dollar seed via PR Newswire: company-reported until an investor confirms independently.",
 "Team and Guardtime history: company-reported.",
 "Headline agent-infrastructure thesis: independently reported.",
 "Engine enrichment on Mike Gault: engine-derived.",
],12.5,gap=8)
txt(s,Inches(6.9),Inches(1.45),Inches(6),Inches(0.4),"Discipline rules",13,NAVY,True)
bullets(s,Inches(6.9),Inches(1.95),Inches(5.9),Inches(4.6),[
 "Company-reported metrics are not treated as verified.",
 "Funding does not prove demand.",
 "GitHub activity does not prove commercial traction.",
 "Token value is not assumed to reach equity.",
 "No public Headline investment or interest in Unicity was identified. That is not the same as confirming none exists.",
 "Full claim-by-claim ledger in AOS_Unicity_Source_Ledger.xlsx.",
],12.5,gap=8)
footer(s,12)

prs.save(os.path.join(F.OUTDIR,"AOS_Unicity_Investment_Deck.pptx"))
print("deck built:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
